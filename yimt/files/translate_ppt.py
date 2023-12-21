""""PPT file translation"""
import argparse

import os

from pptx import Presentation
from yimt.api.utils import detect_lang

ppt_progress = ""
def scan_doc(ppt, new_ppt):
    """Get text to be translated"""
    runs = []
    for i, (slide, new_slide) in enumerate(zip(ppt.slides, new_ppt.slides)):
        # print("Slide{}".format(i + 1), slide.slide_layout)
        for j, (shape, new_shape) in enumerate(zip(slide.shapes, new_slide.shapes)):
            # print("Shape{}".format(j + 1), shape)
            if shape.has_text_frame:
                text_frame = shape.text_frame
                for k, (paragraph, new_paragraph) in enumerate(
                        zip(text_frame.paragraphs, new_shape.text_frame.paragraphs)):
                    # print("\tParagraph{}".format(k + 1), paragraph.text)
                    runs.append(new_paragraph)
    return runs


def translate_ppt_auto(in_fn, source_lang="auto", target_lang="zh", translation_file=None):
    paths = os.path.splitext(in_fn)
    docx_fn = in_fn

    if translation_file is None:
        translated_fn = paths[0] + "-translated.pptx"
    else:
        translated_fn = translation_file

    doc = Presentation(docx_fn)
    translated_doc = Presentation(docx_fn)
    runs = scan_doc(doc, translated_doc)  # get text paragraphs

    if source_lang == "auto":
        source_lang = detect_lang(runs[0].text)

    from yimt.api.translators import Translators
    translator = Translators().get_translator(source_lang, target_lang)

    txt_list = [r.text for r in runs]
    global ppt_progress
    ppt_progress = ""
    batch_size = 10  # 每多少个文本更新一个进度单位
    result_list = []
    for i in range(0, len(txt_list) // batch_size + 1):
        batch = txt_list[i * batch_size: i * batch_size + batch_size]
        result = translator.translate_list(batch)
        result_list += result
        # print(batch) # 测试用
        ppt_progress += "#"
        print("ppt_progress:" + ppt_progress)  # 测试用

    # result_list = translator.translate_list(txt_list)  # translate
    for i in range(len(runs)):
        runs[i].text = result_list[i]

    translated_doc.save(translated_fn)
    ppt_progress = ""  # 每完成一次文件翻译，归零进度
    return translated_fn


if __name__ == "__main__":
    arg_parser = argparse.ArgumentParser("PPT File Translator")
    arg_parser.add_argument("--to_lang", type=str, default="zh", help="target language")
    arg_parser.add_argument("--input_file", type=str, required=True, help="file to be translated")
    arg_parser.add_argument("--output_file", type=str, default=None, help="translation file")
    args = arg_parser.parse_args()

    to_lang = args.to_lang
    in_file = args.input_file
    out_file = args.output_file

    translated_fn = translate_ppt_auto(in_file, target_lang=to_lang, translation_file=out_file)

    import webbrowser

    webbrowser.open(in_file)
    webbrowser.open(translated_fn)

